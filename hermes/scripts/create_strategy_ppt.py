#!/usr/bin/env python3
"""
渔芯科技 — 双业务线战略分析 PPT 生成器
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import datetime

# ── 配色方案 ──
BLUE_DARK = RGBColor(0x1B, 0x2A, 0x4A)   # 深蓝
BLUE_MID = RGBColor(0x2D, 0x5B, 0x8E)    # 中海蓝
BLUE_LIGHT = RGBColor(0x4A, 0x90, 0xD9)   # 浅蓝
TEAL = RGBColor(0x00, 0x96, 0x88)         # 青绿
ORANGE = RGBColor(0xF5, 0x7C, 0x00)       # 橙色
RED = RGBColor(0xE5, 0x3E, 0x3E)          # 红色
GREEN = RGBColor(0x38, 0xA1, 0x69)        # 绿色
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x33, 0x33, 0x33)
GRAY = RGBColor(0x90, 0x90, 0x90)
BG_LIGHT = RGBColor(0xF5, 0xF7, 0xFA)

prs = Presentation()
prs.slide_width = Inches(13.333)   # 16:9 widescreen
prs.slide_height = Inches(7.5)

# ── Helper functions ──

def add_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def add_shape(slide, left, top, width, height, color, shape_type=MSO_SHAPE.RECTANGLE):
    shape = slide.shapes.add_shape(shape_type, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text_box(slide, left, top, width, height, text, font_size=14, color=BLACK, bold=False, alignment=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return txBox

def add_rich_text(slide, left, top, width, height, lines, default_size=12):
    """lines: list of (text, font_size, color, bold)"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        text, size, color, bold = line[0], line[1] if len(line)>1 else default_size, line[2] if len(line)>2 else BLACK, line[3] if len(line)>3 else False
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.color.rgb = color
        p.font.bold = bold
        p.font.name = "Microsoft YaHei"
        p.space_after = Pt(4)
    return txBox

def add_card(slide, left, top, width, height, title, items, title_color=WHITE, bg_color=None, item_color=WHITE):
    if bg_color is None:
        bg_color = BLUE_MID
    card = add_shape(slide, left, top, width, height, bg_color, MSO_SHAPE.ROUNDED_RECTANGLE)
    card.shadow.inherit = False
    add_text_box(slide, left + Inches(0.15), top + Inches(0.08), width - Inches(0.3), Inches(0.4), title, font_size=13, color=title_color, bold=True)
    content = "\n".join(f"• {item}" for item in items)
    add_text_box(slide, left + Inches(0.15), top + Inches(0.55), width - Inches(0.3), height - Inches(0.65), content, font_size=10, color=item_color)
    return card

def slide_title(slide, title, subtitle=""):
    """Add a standard title bar"""
    add_shape(slide, Inches(0), Inches(0), prs.slide_width, Inches(1.1), BLUE_DARK)
    add_text_box(slide, Inches(0.8), Inches(0.15), Inches(11), Inches(0.7), title, font_size=28, color=WHITE, bold=True)
    if subtitle:
        add_text_box(slide, Inches(0.8), Inches(0.7), Inches(11), Inches(0.4), subtitle, font_size=14, color=RGBColor(0xBB,0xCC,0xDD))

# ==================== SLIDE 1: COVER ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
add_bg(slide, BLUE_DARK)

# Decorative line
add_shape(slide, Inches(0), Inches(3.0), prs.slide_width, Inches(0.04), TEAL)
add_shape(slide, Inches(0), Inches(5.5), prs.slide_width, Inches(0.02), TEAL)

add_text_box(slide, Inches(1), Inches(1.2), Inches(11), Inches(1), "渔 芯 科 技", font_size=52, color=WHITE, bold=True, alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(2.2), Inches(11), Inches(0.6), "YUXIN TECHNOLOGY", font_size=18, color=GRAY, alignment=PP_ALIGN.LEFT)
add_text_box(slide, Inches(1), Inches(3.3), Inches(11), Inches(1.2),
    "双业务线战略 · 产品研发布局 · 完整架构分析",
    font_size=24, color=RGBColor(0xBB,0xCC,0xDD), alignment=PP_ALIGN.LEFT)
add_rich_text(slide, Inches(1), Inches(4.5), Inches(11), Inches(0.8), [
    ("渔芯科技1  —  水产养殖之芯（垂直行业深度）", 16, TEAL, True),
    ("渔芯科技2  —  授人以渔通用创业平台（生态孵化广度）", 16, ORANGE, True),
])
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
    f"产品研发中心  |  {datetime.date.today().strftime('%Y年%m月%d日')}",
    font_size=12, color=GRAY)

# ==================== SLIDE 2: 公司战略概览 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_title(slide, "公司战略概览：双业务线并行")

# Left: 渔芯1
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(5.5), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_shape(slide, Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.6), BLUE_MID, MSO_SHAPE.RECTANGLE)
add_text_box(slide, Inches(0.7), Inches(1.55), Inches(5.4), Inches(0.5), "🐟 渔芯科技1 — 水产养殖之芯", font_size=18, color=WHITE, bold=True)
add_rich_text(slide, Inches(0.8), Inches(2.4), Inches(5.2), Inches(4.4), [
    ("定位：RAS循环水养殖全产业链AI自动化", 14, BLUE_DARK, True),
    ("", 6, BLACK),
    ("五环价值链（行业垂直深耕）：", 13, BLUE_MID, True),
    ("  ① 收集资料 — RKR知识库 + FindEra寻元", 12, BLACK),
    ("  ② 开发产品 — EDAI 9大硬件设备AI开发助手", 12, BLACK),
    ("  ③ 验证产品 — Eq-Sim(13仿真) + LookForge(3D+BOM)", 12, BLACK),
    ("  ④ 生产产品 — 供应链管理体系（外包制造）", 12, BLACK),
    ("  ⑤ 销售产品 — 11模块覆盖外贸全链路", 12, BLACK),
    ("", 6, BLACK),
    ("核心产品矩阵（20个模块）：", 13, BLUE_MID, True),
    ("  • 4模块实际运行 | 4模块Docker就绪 | 11模块刚搭建", 12, BLACK),
    ("  • 知识层: RKR(8000) + FindEra(8003)", 12, BLACK),
    ("  • 仿真层: AquaForge + Eq-Sim + LookForge", 12, BLACK),
    ("  • 硬件层: EDAI 9子项目(各自独立进化)", 12, BLACK),
    ("  • 销售层: MatrixFlow + CRM + 外贸11件套", 12, BLACK),
])

# Right: 渔芯2
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(5.5), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
add_shape(slide, Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.6), ORANGE, MSO_SHAPE.RECTANGLE)
add_text_box(slide, Inches(7.0), Inches(1.55), Inches(5.4), Inches(0.5), "🎓 渔芯科技2 — 授人以渔通用平台", font_size=18, color=WHITE, bold=True)
add_rich_text(slide, Inches(7.1), Inches(2.4), Inches(5.2), Inches(4.4), [
    ("定位：AI驱动的产品全生命周期通用操作系统", 14, BLUE_DARK, True),
    ("", 6, BLACK),
    ("核心资产（从渔芯1蒸馏）：", 13, ORANGE, True),
    ("  • 通用知识中枢 — RKR可适配任意行业", 12, BLACK),
    ("  • 自进化框架SDK — 任意模块可接入知识进化", 12, BLACK),
    ("  • 设备开发框架 — EDAI方法论通用化", 12, BLACK),
    ("  • 仿真验证平台 — 参数化仿真引擎", 12, BLACK),
    ("  • 销售电商中台 — 外贸全链路可复用", 12, BLACK),
    ("  • 07-DevPlan — AI一键生成完整开发文档", 12, BLACK),
    ("", 6, BLACK),
    ("生态社区（OPC通用管理平台）：", 13, ORANGE, True),
    ('  🔧"我开发" — 专注产品研发，使用AI开发工具', 12, BLACK),
    ('  🏭"你生产" — 专注生产制造，接单生产', 12, BLACK),
    ('  📦"他销售" — 专注市场销售，共享产品红利', 12, BLACK),
    ("", 4, BLACK),
    ("愿景：降低创业门槛，共享项目红利", 13, RED, True),
])

# ==================== SLIDE 3: 五环价值链 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_title(slide, "五环价值链 × 双业务线矩阵")

rings = [
    ("① 收集资料", "RKR知识库\nFindEra寻元", "通用知识中枢\n行业无关调研引擎", BLUE_MID),
    ("② 开发产品", "EDAI 9设备\nAquaForge设计", "通用设备开发框架\nAI辅助产品研发", TEAL),
    ("③ 验证产品", "Eq-Sim 13仿真\nLookForge 3D+BOM", "通用仿真验证平台\n参数化仿真引擎", GREEN),
    ("④ 生产产品", "供应链管理\n(外包制造)", "供应链协作网络\n生产外包管理", ORANGE),
    ("⑤ 销售产品", "外贸11模块\nMatrixFlow矩阵", "销售电商中台\n全链路可复用", RED),
]

x_start = 0.3
ring_w = 2.3
gap = 0.25
y_ring = 1.5
ring_h = 3.2
y_label = 5.0

for i, (title, left_text, right_text, color) in enumerate(rings):
    x = Inches(x_start + i * (ring_w + gap))
    # Ring card
    card = add_shape(slide, x, Inches(y_ring), Inches(ring_w), Inches(ring_h), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    # Color header
    add_shape(slide, x, Inches(y_ring), Inches(ring_w), Inches(0.55), color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.1), Inches(y_ring + 0.05), Inches(ring_w - 0.2), Inches(0.45), title, font_size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    # Content
    add_text_box(slide, x + Inches(0.15), Inches(y_ring + 0.7), Inches(ring_w - 0.3), Inches(1.2),
        f"渔芯1:\n{left_text}", font_size=10, color=BLACK)
    add_shape(slide, x + Inches(0.15), Inches(y_ring + 2.0), Inches(ring_w - 0.3), Inches(0.015), GRAY)
    add_text_box(slide, x + Inches(0.15), Inches(y_ring + 2.15), Inches(ring_w - 0.3), Inches(1.0),
        f"渔芯2:\n{right_text}", font_size=10, color=color)
    # Arrow between rings
    if i < 4:
        add_text_box(slide, x + Inches(ring_w), Inches(y_ring + 1.3), Inches(gap), Inches(0.5), "→", font_size=22, color=GRAY, bold=True, alignment=PP_ALIGN.CENTER)

# Bottom: flow
add_text_box(slide, Inches(1), Inches(6.2), Inches(11), Inches(0.5),
    "五环闭环: 知识驱动产品研发 → 仿真验证设计方案 → 外包生产制造 → 全链路销售交付 → 数据反馈优化知识",
    font_size=12, color=BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# ==================== SLIDE 4: 产品研发全景图 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_title(slide, "产品研发全景：25个模块按价值链分层", "6-产品研发/ 完整目录结构")

# Layer definitions
layers = [
    ("基础设施层 (3模块)", ["00-FindEra寻元 ✅", "01-RKR知识库 ✅", "00-公共组件 📋"], BLUE_DARK, 0.8),
    ("养殖仿真层 (4模块)", ["02-AquaForge 📦", "03-EDAI(9子) v1.0", "04-Eq-Sim 📦", "05-LookForge 📦"], BLUE_MID, 1.0),
    ("软件开发层 (1模块)", ["07-DevPlan v2.0 ✅"], TEAL, 0.5),
    ("销售运营层 (13模块)", ["08-MatrixFlow ✅", "08-19运营面板 📋", "09~19外贸11件套 🆕"], ORANGE, 1.0),
    ("平台层 (2模块)", ["OPC通用管理平台 📋", "历史版本(5+) 🗄️"], GRAY, 0.6),
]

y = 1.45
for title, modules, color, height in layers:
    add_shape(slide, Inches(0.4), Inches(y), Inches(2.8), Inches(height), color, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_text_box(slide, Inches(0.55), Inches(y + 0.05), Inches(2.5), Inches(height - 0.1), title, font_size=12, color=WHITE, bold=True)
    content = "\n".join(f"  {m}" for m in modules)
    add_text_box(slide, Inches(3.5), Inches(y + 0.05), Inches(9.3), Inches(height - 0.1), content, font_size=11, color=BLACK)
    y += height + 0.12

# Legend
add_rich_text(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.8), [
    ("图例: ✅运行中(4)  📦Docker就绪(4)  🆕刚搭建(11)  📋规划中(3)  🗄️已归档(5+)", 11, GRAY),
    ("技术栈统一: FastAPI + 独立SQLite + Docker独立部署 + version.json + self_evolution", 11, GRAY),
])

# ==================== SLIDE 5: EDAI 9子项目 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_title(slide, "EDAI 硬件开发 — 9大RAS设备AI开发助手", "各自独立进化 + 统一自进化框架 + VersionManager")

edai_items = [
    ("增氧曝气系统", "曝气量/气泡尺寸/氧传递效率", "🔵"),
    ("滚筒微滤机", "过滤精度/堵塞周期/反冲洗优化", "🔵"),
    ("生物过滤设备", "MBBR填料/硝化速率/氨氮去除", "🟢"),
    ("蛋白分离器", "气浮效率/有机物去除/臭氧耦合", "🟢"),
    ("温控设备", "加热制冷功率/保温材料/能耗分析", "🟠"),
    ("消毒设备", "UV剂量/臭氧浓度/接触时间", "🟠"),
    ("智能控制", "PLC逻辑/传感器选型/控制算法", "🔴"),
    ("自动投喂设备", "投喂策略/生长模型/饲料转化率", "🔴"),
    ("数据采集设备", "水质传感器/采样频率/数据清洗", "🟣"),
]

col_w = 4.0
row_h = 1.5
for idx, (name, desc, icon) in enumerate(edai_items):
    col = idx % 3
    row = idx // 3
    x = Inches(0.5 + col * col_w)
    y = Inches(1.45 + row * row_h)
    add_shape(slide, x, y, Inches(col_w - 0.3), Inches(row_h - 0.15), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x, y, Inches(0.08), Inches(row_h - 0.15), BLUE_MID, MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.1), Inches(col_w - 0.6), Inches(0.45),
        f"{icon} {name}", font_size=14, color=BLUE_DARK, bold=True)
    add_text_box(slide, x + Inches(0.2), y + Inches(0.6), Inches(col_w - 0.6), Inches(0.5),
        f"核心计算: {desc}", font_size=10, color=GRAY)

# Note
add_rich_text(slide, Inches(0.5), Inches(6.2), Inches(12), Inches(0.8), [
    ("每个子项目 = FastAPI独立后端 + 独立evolution.py + 独立VersionManager + 独立Docker部署", 12, BLUE_MID, True),
    ("共享: self_evolution SDK / version_manager / research_request_sdk / 统一端口规范(8101-8109,3101-3109)", 11, GRAY),
    ("所有设备研发模块均通过自进化框架定期扫描RKR，获取领域最新知识，发现缺口自动申请调研", 11, GRAY),
])

# ==================== SLIDE 6: 仿真验证层 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_title(slide, "仿真验证层：三件套覆盖从设备到系统的完整验证链")

sims = [
    ("Eq-Sim 设备仿真", "13种专业仿真",
     ["流场仿真 | 溶氧分布 | 温度场 | 管网水力",
      "生物滤池(Ergun) | MBBR(Monod) | 沉淀池(Stokes)",
      "氧气锥(K_La) | 蛋白分离器 | 滚筒微滤机",
      "结构强度 | 热仿真 | ROI回本分析"],
      GREEN, Inches(0.5)),
    ("LookForge 系统设计", "3D布局 + 管道路由 + BOM",
     ["3D布局设计器(32种设备GLB渲染)",
      "管道自动路由(A*转角最小路径)",
      "BOM物料清单(设备+管材+管件+费用+能耗)",
      "水力计算器(23种工程公式)", "物种参数库(50+品种)"],
      TEAL, Inches(4.5)),
    ("AquaForge 数字孪生", "4模型联动全场景仿真",
     ["天气模型 → 水质模型 → 生长模型 → 设备模型",
      "AI决策规则引擎(投喂/设备/告警)",
      "模板选择式养殖方案", "游戏化激励系统(成就/签到/排行)",
      "管理后台(运营数据面板)"],
      BLUE_MID, Inches(8.5)),
]

for title, subtitle, items, color, x_pos in sims:
    add_shape(slide, x_pos, Inches(1.45), Inches(3.7), Inches(5.2), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x_pos, Inches(1.45), Inches(3.7), Inches(0.55), color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x_pos + Inches(0.15), Inches(1.5), Inches(3.4), Inches(0.45), title, font_size=15, color=WHITE, bold=True)
    add_text_box(slide, x_pos + Inches(0.15), Inches(2.1), Inches(3.4), Inches(0.3), subtitle, font_size=11, color=color, bold=True)
    content = "\n".join(f"  ✓ {item}" for item in items)
    add_text_box(slide, x_pos + Inches(0.15), Inches(2.5), Inches(3.4), Inches(3.8), content, font_size=11, color=BLACK)

add_text_box(slide, Inches(0.5), Inches(6.85), Inches(12), Inches(0.4),
    "验证链: EDAI设备设计方案 → Eq-Sim单设备性能验证 → LookForge系统集成验证 → AquaForge全场景数字孪生",
    font_size=13, color=BLUE_DARK, bold=True, alignment=PP_ALIGN.CENTER)

# ==================== SLIDE 7: 销售外贸层 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_title(slide, "销售外贸层：11模块覆盖外贸全链路", "统一架构: FastAPI + SQLite + Docker Compose + version.json")

sales_items = [
    ("09-线索挖掘", "获客", "线索发现+评分"), ("10-企微营销", "触达", "私域运营"),
    ("11-RAS报价引擎", "报价", "BOM+PDF生成"), ("17-智能CRM", "管理", "客户/商机/漏斗"),
    ("12-合同管理", "成交", "合同生成+签署"), ("15-外贸跟单", "履约", "订单进度跟踪"),
    ("16-国际物流", "交付", "物流追踪"), ("13-外贸询盘", "外贸", "询盘处理"),
    ("14-多语言产品站", "展示", "国际化展示"), ("18-竞品情报", "情报", "竞品监控"),
    ("19-销售知识库", "赋能", "销售资料管理"), ("08-MatrixFlow", "内容", "自媒体矩阵 ✅"),
]

for idx, (name, stage, desc) in enumerate(sales_items):
    col = idx % 4
    row = idx // 4
    x = Inches(0.5 + col * 3.1)
    y = Inches(1.45 + row * 1.7)
    row_color = [BLUE_MID, TEAL, GREEN, ORANGE, RED][row]
    add_shape(slide, x, y, Inches(2.8), Inches(1.5), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x, y, Inches(2.8), Inches(0.4), row_color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.02), Inches(2.6), Inches(0.35), f"{name} [{stage}]", font_size=11, color=WHITE, bold=True)
    add_text_box(slide, x + Inches(0.1), y + Inches(0.55), Inches(2.6), Inches(0.7), desc, font_size=10, color=BLACK)

add_rich_text(slide, Inches(0.5), Inches(6.3), Inches(12), Inches(0.9), [
    ("全链路: 获客→触达→报价→管理→成交→履约→交付  |  全部v1.0.0 (2026-06-06创建)", 12, BLUE_MID, True),
    ("仅 MatrixFlow 实际运行，其余10个模块为脚手架模板，待填充业务逻辑", 11, GRAY),
])

# ==================== SLIDE 8: OPC 通用平台 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_title(slide, "渔芯科技2 — OPC 通用管理平台", "创业孵化器 · 生态社区基础 · 授人以渔")

opc_modules = [
    ("1-项目调研", "AI自动调研赛道", "对接FindEra", BLUE_MID),
    ("2-知识库", "建立行业知识体系", "对接RKR", BLUE_LIGHT),
    ("3-竞品分析", "AI分析竞争对手", "对接FindEra", TEAL),
    ("4-开发", "AI辅助产品开发", "对接07+EDAI", GREEN),
    ("5-供应链&材料库", "对接材料供应商", "外包生产管理", ORANGE),
    ("6-销售", "一键接入销售中台", "对接09-19", RED),
    ("7-客服", "AI客服系统", "对接Hermes", RGBColor(0x9C,0x27,0xB0)),
]

for idx, (name, desc, hook, color) in enumerate(opc_modules):
    x = Inches(0.3 + idx * 1.8)
    # Card
    add_shape(slide, x, Inches(1.6), Inches(1.65), Inches(3.8), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x, Inches(1.6), Inches(1.65), Inches(1.65), color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.08), Inches(1.65), Inches(1.5), Inches(0.5), name, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.08), Inches(2.1), Inches(1.5), Inches(0.35), "━━━━", font_size=8, color=WHITE, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.08), Inches(2.3), Inches(1.5), Inches(0.8), desc, font_size=11, color=BLACK, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.08), Inches(3.8), Inches(1.5), Inches(0.4), f"→ {hook}", font_size=9, color=color, alignment=PP_ALIGN.CENTER)
    # Arrow
    if idx < 6:
        add_text_box(slide, x + Inches(1.65), Inches(2.8), Inches(0.2), Inches(0.4), "→", font_size=16, color=GRAY, bold=True)

# Bottom: eco vision
add_rich_text(slide, Inches(0.5), Inches(5.7), Inches(12), Inches(1.5), [
    ("生态社区愿景", 18, BLUE_DARK, True),
    ("", 6, BLACK),
    ('🔧 "我开发" — 专注产品研发，使用AI开发工具（EDAI + 07-DevPlan + Eq-Sim）', 13, BLACK),
    ('🏭 "你生产" — 专注生产制造，接入供应链网络（OPC-5 供应链&材料库）', 13, BLACK),
    ('📦 "他销售" — 专注市场销售，共享项目红利（09-19 销售矩阵）', 13, BLACK),
    ("", 4, BLACK),
    ("核心创新：无需打通全部五环，只需专精一环即可通过OPC平台参与项目并共享红利", 13, ORANGE, True),
])

# ==================== SLIDE 9: 战略路线图 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BG_LIGHT)
slide_title(slide, "战略路线图：四阶段推进", "从夯实基础到生态社区的完整路径")

phases = [
    ("Phase 1\n(当前)", "夯实知识层\n激活仿真核心",
     ["✅ RKR + FindEra 运行中", "🔧 AquaForge 激活", "🔧 Eq-Sim 激活", "🔧 LookForge 激活", "🔧 自进化框架全模块推广"],
     BLUE_MID),
    ("Phase 2\n(短期)", "蒸馏通用平台\n填充开发工具",
     ["🔧 OPC-1~4 填充实际功能", "🔧 07-DevPlan 接入OPC-4", "🔧 版本管理器规范化", "🔧 调研请求规范落地", "🔧 文档质量自迭代系统上线"],
     TEAL),
    ("Phase 3\n(中期)", "销售线充实\n供应链对接",
     ["🔧 09-19 从模板到可用", "🔧 OPC-5 供应链实现", "🔧 外贸全链路打通", "🔧 MatrixFlow 策略引擎", "🔧 多语言产品站上线"],
     ORANGE),
    ("Phase 4\n(长期)", "生态社区\n开放平台",
     ["🔧 多租户 + 角色权限", "🔧 项目红利分配机制", "🔧 社区市场(模块交易)", "🔧 服务撮合引擎", "🔧 公开SDK + 文档站"],
     RED),
]

for idx, (title, subtitle, items, color) in enumerate(phases):
    x = Inches(0.4 + idx * 3.2)
    # Phase card
    add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(5.3), WHITE, MSO_SHAPE.ROUNDED_RECTANGLE)
    add_shape(slide, x, Inches(1.5), Inches(2.9), Inches(1.2), color, MSO_SHAPE.RECTANGLE)
    add_text_box(slide, x + Inches(0.15), Inches(1.55), Inches(2.6), Inches(0.7), title, font_size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(slide, x + Inches(0.15), Inches(2.3), Inches(2.6), Inches(0.3), subtitle, font_size=12, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    content = "\n".join(items)
    add_text_box(slide, x + Inches(0.15), Inches(2.8), Inches(2.6), Inches(3.5), content, font_size=10, color=BLACK)
    # Connector
    if idx < 3:
        add_text_box(slide, x + Inches(2.9), Inches(3.5), Inches(0.35), Inches(0.4), "▸", font_size=24, color=GRAY, bold=True)

# ==================== SLIDE 10: 总结 ====================
slide = prs.slides.add_slide(prs.slide_layouts[6])
add_bg(slide, BLUE_DARK)

add_shape(slide, Inches(0), Inches(3.2), prs.slide_width, Inches(0.03), TEAL)

add_text_box(slide, Inches(1), Inches(0.8), Inches(11), Inches(0.8), "总结与愿景", font_size=36, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(slide, Inches(1), Inches(1.6), Inches(11), Inches(0.6),
    "渔芯科技 = 水产养殖全产业链AI自动化 × 通用创业平台",
    font_size=18, color=TEAL, alignment=PP_ALIGN.CENTER)

key_points = [
    "✅ 五环价值链完整覆盖：收集资料 → 开发产品 → 验证产品 → 生产产品 → 销售产品",
    "✅ 20个产品模块，统一架构规范：FastAPI + SQLite + Docker + 自进化框架",
    "✅ 双业务线并行：垂直行业深耕（渔芯1）+ 通用平台孵化（渔芯2）",
    "✅ 核心技术壁垒：RAS仿真三件套 + EDAI 9设备 + 13种工程仿真",
    "✅ 知识驱动自进化：RKR知识中枢 → 自进化SDK → 所有模块持续进化",
    "🎯 战略终局：OPC通用管理平台 → \"我开发，你生产，他销售\"生态社区",
]

for i, point in enumerate(key_points):
    add_text_box(slide, Inches(1.5), Inches(3.5 + i * 0.55), Inches(10), Inches(0.5),
        point, font_size=14, color=WHITE)

add_text_box(slide, Inches(1), Inches(6.8), Inches(11), Inches(0.4),
    "渔芯科技 — 授人以鱼，不如授人以渔",
    font_size=16, color=GRAY, alignment=PP_ALIGN.CENTER)

# ── Save ──
output_path = "/Users/hua/Desktop/渔芯科技_双业务线战略分析.pptx"
prs.save(output_path)
print(f"PPT saved to: {output_path}")
print(f"Slides: {len(prs.slides)}")
